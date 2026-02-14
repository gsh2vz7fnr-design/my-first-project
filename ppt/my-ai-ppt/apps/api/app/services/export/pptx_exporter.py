from __future__ import annotations

from pathlib import Path
from typing import Dict, Iterable, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.schemas.dsl import ProjectDSL, Slide
from app.services.export.layout_slots import SLOT_MAP

SLIDE_W = 13.333
SLIDE_H = 7.5
EXPORT_ROOT = Path("/tmp/my-ai-ppt-exports")


def _pct_to_inches(value: str, total_inches: float) -> float:
    return (float(value.replace("%", "")) / 100.0) * total_inches


def _slot_box(slot: Dict[str, str]) -> Tuple[float, float, float, float]:
    x = _pct_to_inches(slot["x"], SLIDE_W)
    y = _pct_to_inches(slot["y"], SLIDE_H)
    w = _pct_to_inches(slot["w"], SLIDE_W)
    h = _pct_to_inches(slot["h"], SLIDE_H)
    return x, y, w, h


def _add_textbox(
    ppt_slide,
    text: str,
    slot: Dict[str, str],
    font_size: int,
    bold: bool = False,
    align: str = "left",
    color: Tuple[int, int, int] = (17, 24, 39),
) -> None:
    if not text:
        return
    x, y, w, h = _slot_box(slot)
    shape = ppt_slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    para = frame.paragraphs[0]
    para.text = text
    para.font.size = Pt(font_size)
    para.font.bold = bold
    para.font.color.rgb = RGBColor(*color)
    if align == "center":
        para.alignment = PP_ALIGN.CENTER
    elif align == "right":
        para.alignment = PP_ALIGN.RIGHT
    else:
        para.alignment = PP_ALIGN.LEFT


def _add_bullets(ppt_slide, items: Iterable[str], slot: Dict[str, str], font_size: int = 20) -> None:
    bullets = [item for item in items if item]
    if not bullets:
        return
    x, y, w, h = _slot_box(slot)
    shape = ppt_slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    for index, item in enumerate(bullets):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.text = item
        para.level = 0
        para.font.size = Pt(font_size)
        para.font.color.rgb = RGBColor(15, 23, 42)
        para.font.bold = False


def _render_cover(ppt_slide, dsl_slide: Slide) -> None:
    slots = SLOT_MAP["cover_centered_01"]
    _add_textbox(ppt_slide, dsl_slide.content.title or "", slots["title"], 42, bold=True, align="center")
    _add_textbox(ppt_slide, dsl_slide.content.subtitle or "", slots["subtitle"], 24, align="center", color=(51, 65, 85))
    _add_textbox(ppt_slide, dsl_slide.content.footer or "", slots["footer"], 13, align="center", color=(100, 116, 139))


def _render_split(ppt_slide, dsl_slide: Slide) -> None:
    slots = SLOT_MAP["split_left_image_right_text"]
    _add_textbox(ppt_slide, dsl_slide.content.title or "", slots["title"], 34, bold=True)
    _add_textbox(ppt_slide, dsl_slide.content.subtitle or "", slots["subtitle"], 19, color=(51, 65, 85))
    _add_bullets(ppt_slide, dsl_slide.content.bullets[:6], slots["bullets"], font_size=19)
    _add_textbox(ppt_slide, dsl_slide.content.footer or "", slots["footer"], 12, color=(100, 116, 139))

    # V1 placeholder for image slot when no image source provided.
    x, y, w, h = _slot_box(slots["imageSlot"])
    rect = ppt_slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))  # 1 == Rectangle
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(239, 246, 255)
    rect.line.color.rgb = RGBColor(191, 219, 254)
    frame = rect.text_frame
    frame.text = "image slot"
    frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    frame.paragraphs[0].font.size = Pt(14)
    frame.paragraphs[0].font.color.rgb = RGBColor(71, 85, 105)


def _render_three_column(ppt_slide, dsl_slide: Slide) -> None:
    slots = SLOT_MAP["three_column_points"]
    _add_textbox(ppt_slide, dsl_slide.content.title or "", slots["title"], 34, bold=True)
    points = (dsl_slide.content.bullets or [])[:3]
    for idx, key in enumerate(["card1", "card2", "card3"]):
        if idx < len(points):
            _add_textbox(ppt_slide, points[idx], slots[key], 20, align="left")
    _add_textbox(ppt_slide, dsl_slide.content.footer or "", slots["footer"], 12, color=(100, 116, 139))


def _render_slide(ppt_slide, dsl_slide: Slide) -> None:
    if dsl_slide.layout_id == "split_left_image_right_text":
        _render_split(ppt_slide, dsl_slide)
        return
    if dsl_slide.layout_id == "three_column_points":
        _render_three_column(ppt_slide, dsl_slide)
        return
    _render_cover(ppt_slide, dsl_slide)


def export_editable_pptx(project_id: str, job_id: str, dsl: ProjectDSL) -> str:
    out_dir = EXPORT_ROOT / project_id
    out_dir.mkdir(parents=True, exist_ok=True)
    out_file = out_dir / f"{job_id}.pptx"

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank_layout = prs.slide_layouts[6]

    for dsl_slide in dsl.slides:
        slide = prs.slides.add_slide(blank_layout)
        _render_slide(slide, dsl_slide)

    prs.save(str(out_file))
    return str(out_file)


def export_image_fallback_pptx(project_id: str, job_id: str, dsl: ProjectDSL) -> str:
    # V1: fall back to editable pipeline until image render worker is available.
    return export_editable_pptx(project_id, f"{job_id}-imgfb", dsl)
