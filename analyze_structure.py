from pptx import Presentation
import json
import re

def extract_content(file_path):
    prs = Presentation(file_path)
    slides_data = []
    
    for i, slide in enumerate(prs.slides):
        slide_info = {
            "index": i,
            "title": "",
            "body": [],
            "layout_suggestion": "generic"
        }
        
        # Sort shapes by position (top to bottom, left to right)
        shapes = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                shapes.append((shape.top, shape.left, shape.text.strip()))
        
        shapes.sort(key=lambda x: (x[0], x[1])) # Sort by Top then Left
        
        if not shapes:
            slides_data.append(slide_info)
            continue
            
        # Assume first item is title if it's near top
        if shapes[0][0] < 2000000: # heuristic for top area
            slide_info["title"] = shapes[0][2]
            body_shapes = shapes[1:]
        else:
            body_shapes = shapes
            
        slide_info["body"] = [s[2] for s in body_shapes]
        
        # Simple heuristics for layout
        if i == 0:
            slide_info["layout_suggestion"] = "cover"
        elif len(slide_info["body"]) >= 3 and any(len(txt) < 50 for txt in slide_info["body"]):
             # Many short items -> Grid or List
            slide_info["layout_suggestion"] = "grid"
        elif "vs" in slide_info["title"].lower() or "对比" in slide_info["title"]:
            slide_info["layout_suggestion"] = "split"
        elif len(slide_info["body"]) == 2:
            slide_info["layout_suggestion"] = "split"
        
        slides_data.append(slide_info)
        
    return slides_data

if __name__ == "__main__":
    data = extract_content("产品经理AI培训 - 彭鹏.pptx")
    print(json.dumps(data, ensure_ascii=False, indent=2))
