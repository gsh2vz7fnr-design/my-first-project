from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import copy

# Configuration
INPUT_FILE = "产品经理AI培训 - 彭鹏.pptx"
OUTPUT_FILE = "产品经理AI培训_深度美化版.pptx"
ASSET_DIR = "assets"

# Colors
COLOR_PRIMARY = RGBColor(0, 51, 102) # Dark Blue
COLOR_ACCENT = RGBColor(0, 102, 204) # Bright Blue
COLOR_BG = RGBColor(245, 247, 250) # Light Grey
COLOR_TEXT_MAIN = RGBColor(60, 60, 60)
COLOR_TEXT_LIGHT = RGBColor(255, 255, 255)
COLOR_CARD_BG = RGBColor(255, 255, 255)

def get_asset(name):
    path = os.path.join(ASSET_DIR, f"asset_{name}.png")
    return path if os.path.exists(path) else None

def clear_slide(slide):
    # Remove all shapes to rebuild
    # iterating copy of list because we are modifying it
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

def set_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_title(slide, text):
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1))
    tf = title_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # Add decorative line
    line = slide.shapes.add_connector(
        1, Inches(0.5), Inches(1.3), Inches(12.83), Inches(1.3)
    )
    line.line.color.rgb = COLOR_ACCENT
    line.line.width = Pt(2)

def create_card(slide, left, top, width, height, title, content=None, icon=None):
    # Card Background (Rounded Rectangle)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_CARD_BG
    shape.line.color.rgb = RGBColor(220, 220, 220)
    shape.shadow.inherit = False # No default shadow, maybe add soft one later if possible
    
    # Icon
    current_top = top + Inches(0.2)
    if icon:
        icon_path = get_asset(icon)
        if icon_path:
            slide.shapes.add_picture(icon_path, left + Inches(0.2), current_top, width=Inches(0.8))
            # If icon exists, title moves right
            text_left = left + Inches(1.1)
            text_width = width - Inches(1.3)
        else:
            text_left = left + Inches(0.2)
            text_width = width - Inches(0.4)
    else:
        text_left = left + Inches(0.2)
        text_width = width - Inches(0.4)

    # Title
    if title:
        tb = slide.shapes.add_textbox(text_left, current_top, text_width, Inches(0.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        current_top += Inches(0.6)

    # Content
    if content:
        # Use simple text box for content inside card
        tb = slide.shapes.add_textbox(left + Inches(0.2), current_top, width - Inches(0.4), height - (current_top - top) - Inches(0.2))
        tf = tb.text_frame
        tf.word_wrap = True
        
        # If content is string, split by newline or simple wrap
        # If content is list, add bullets
        lines = content if isinstance(content, list) else content.split('\n')
        
        for line in lines:
            if not line.strip(): continue
            p = tf.add_paragraph()
            p.text = line.strip()
            p.font.name = 'Microsoft YaHei'
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_TEXT_MAIN
            p.space_after = Pt(6)
            if len(line) > 50: # Long text
                p.alignment = PP_ALIGN.JUSTIFY

def analyze_and_rebuild(prs):
    # We iterate and REBUILD each slide.
    # Note: Deleting shapes and adding new ones is safer than modifying complex existing groups.
    
    # 1. Cover Slide
    slide = prs.slides[0]
    # Extract text first
    texts = [s.text for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
    title = texts[0] if texts else "AI赋能产品经理"
    subtitle = texts[1] if len(texts) > 1 else ""
    
    clear_slide(slide)
    
    # Add Full Background Image
    bg_path = get_asset("bg_cover")
    if bg_path:
        slide.shapes.add_picture(bg_path, 0, 0, prs.slide_width, prs.slide_height)
    
    # Add Title Overlay
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)
    
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_LIGHT
    
    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(24)
        p.font.color.rgb = COLOR_TEXT_LIGHT
        p.space_before = Pt(20)

    # Process other slides
    for i in range(1, len(prs.slides)):
        slide = prs.slides[i]
        
        # Extract content
        shapes = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                shapes.append((shape.top, shape.left, shape.text.strip()))
        shapes.sort(key=lambda x: (x[0], x[1]))
        
        if not shapes: continue
        
        title = shapes[0][2]
        body_items = [s[2] for s in shapes[1:]]
        
        clear_slide(slide)
        set_background(slide)
        add_title(slide, title)
        
        # Determine Layout based on slide index (hardcoded logic for this specific PPT structure)
        # We can also infer from body_items length
        
        # Slide 2: 4 Core Works (Quadrant)
        if i == 1:
            # Grid 2x2
            # Map items to 4 quadrants
            # Items seem to be grouped in original PPT. 
            # Let's try to just distribute the body items into 4 cards if count is high, or just list them.
            # Analyze extracting shows groups: "市场洞察...", "需求分析...", "PRD...", "数据分析..."
            # Let's create 4 big cards
            
            card_width = Inches(5.5)
            card_height = Inches(2.5)
            
            # Card 1 (Top Left)
            create_card(slide, Inches(0.8), Inches(1.8), card_width, card_height, "市场洞察与用户研究", 
                       content=["海量行业调研", "用户原始反馈处理", "用户画像构建"], icon="icon_research")
            
            # Card 2 (Top Right)
            create_card(slide, Inches(6.8), Inches(1.8), card_width, card_height, "需求分析与流程设计", 
                       content=["需求分析与拆解", "业务流程与逻辑查漏", "产品规划"], icon="icon_plan")
                       
            # Card 3 (Bottom Left)
            create_card(slide, Inches(0.8), Inches(4.5), card_width, card_height, "PRD文档与原型设计", 
                       content=["PRD文档编写", "产品原型设计"], icon="icon_structure")
                       
            # Card 4 (Bottom Right)
            create_card(slide, Inches(6.8), Inches(4.5), card_width, card_height, "数据分析和项目汇报", 
                       content=["数据分析", "项目汇报"], icon="icon_data")

        # Slide 3: Prompt AI vs Skills (Split)
        elif i == 2:
            # Left: Prompt AI
            create_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5), "Prompt AI 提示词工程", 
                       content="连接人类需求与AI能力的桥梁...\n工具: DeepSeek, Gemini", icon="icon_brain")
            
            # Right: Skills
            create_card(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(5), "Skills 技能/智能体", 
                       content="你的数字化分身，入职即巅峰...\n工具: Trae, Custom Gems", icon="icon_robot")

        # Slide 6: Prompt Structures (4 Columns)
        elif i == 5:
            col_width = Inches(2.8)
            spacing = Inches(0.2)
            start_left = Inches(0.5)
            
            titles = ["ICIO 架构", "CO-STAR 架构", "TRDF 架构", "思维链 CoT"]
            descs = ["数据处理与文档生成", "战略思考与汇报沟通", "高频碎片化任务", "复杂逻辑推演"]
            
            for idx in range(4):
                create_card(slide, start_left + (col_width + spacing) * idx, Inches(2), col_width, Inches(4), 
                           titles[idx], content=descs[idx], icon="icon_structure")

        # Slide 7: Skills 4 Elements (Process/Flow)
        elif i == 6:
            # 4 Horizontal Cards with arrows implies flow
            card_width = Inches(2.5)
            spacing = Inches(0.5)
            start_left = Inches(0.8)
            top = Inches(3)
            
            steps = [
                ("定义身份", "Role - 灵魂"),
                ("喂入知识", "Knowledge - 业务Wiki"),
                ("赋予工具", "Capability - 功能插件"),
                ("编排流程", "Workflow - 业务流程图")
            ]
            
            for idx, (t, c) in enumerate(steps):
                create_card(slide, start_left + (card_width + spacing) * idx, top, card_width, Inches(3), 
                           t, content=c, icon="icon_process")
                
                # Add Arrow if not last
                if idx < 3:
                    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, 
                                          start_left + (card_width + spacing) * idx + card_width + Inches(0.1), 
                                          top + Inches(1.2), Inches(0.3), Inches(0.3))
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = COLOR_ACCENT
                    arrow.line.fill.background()

        # Default: List view or Grid for others
        else:
            # Just put body items in a nice large card or two columns
            if len(body_items) > 6:
                # Two columns
                mid = len(body_items) // 2
                create_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5), "Content Part 1", content=body_items[:mid])
                create_card(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(5), "Content Part 2", content=body_items[mid:])
            else:
                create_card(slide, Inches(1.5), Inches(2), Inches(10), Inches(4.5), "", content=body_items)

def main():
    if not os.path.exists(INPUT_FILE):
        print("Input file missing")
        return
        
    prs = Presentation(INPUT_FILE)
    
    # Resize to 16:9 just in case
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    analyze_and_rebuild(prs)
    
    prs.save(OUTPUT_FILE)
    print(f"Remastered presentation saved to {OUTPUT_FILE}")

if __name__ == "__main__":
    main()
